import os
import io
import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


def pdf_to_pptx(pdf_path, pptx_path, dpi=200):
    """
    PDF dosyasını PPTX'e çevirir.
    Her PDF sayfasını yüksek çözünürlüklü görsel olarak alır
    ve PowerPoint'te bir slide'a tam sığdırır.

    Parametreler:
        pdf_path (str): Giriş PDF dosya yolu
        pptx_path (str): Çıkış PPTX dosya yolu
        dpi (int): Render kalitesi. 150-250 arası genelde idealdir.
    """

    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF bulunamadı: {pdf_path}")

    # Yeni sunum oluştur
    prs = Presentation()

    # 16:9 slide boyutu
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # PDF aç
    doc = fitz.open(pdf_path)

    # DPI -> zoom dönüşümü
    zoom = dpi / 72.0
    matrix = fitz.Matrix(zoom, zoom)

    for page_index in range(len(doc)):
        page = doc.load_page(page_index)

        # Sayfayı render et
        pix = page.get_pixmap(matrix=matrix, alpha=False)

        # Pixmap -> PIL Image
        img_bytes = pix.tobytes("png")
        pil_img = Image.open(io.BytesIO(img_bytes))

        # Geçici bellek dosyası yerine bytes kullanacağız
        img_stream = io.BytesIO()
        pil_img.save(img_stream, format="PNG")
        img_stream.seek(0)

        # Yeni boş slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Slide ve görsel oranları
        slide_w = prs.slide_width
        slide_h = prs.slide_height

        img_w_px, img_h_px = pil_img.size
        img_ratio = img_w_px / img_h_px
        slide_ratio = slide_w / slide_h

        # Oranı koruyarak slide içine yerleştir
        if img_ratio > slide_ratio:
            # Görsel daha geniş
            pic_w = slide_w
            pic_h = int(slide_w / img_ratio)
            left = 0
            top = int((slide_h - pic_h) / 2)
        else:
            # Görsel daha uzun
            pic_h = slide_h
            pic_w = int(slide_h * img_ratio)
            top = 0
            left = int((slide_w - pic_w) / 2)

        slide.shapes.add_picture(img_stream, left, top, width=pic_w, height=pic_h)

    # Eğer Presentation() oluşturunca boş ilk slide kalmasın istiyorsan:
    # Varsayılan olarak boş slide gelmez, çünkü add_slide ile ekliyoruz.

    prs.save(pptx_path)
    doc.close()

    print(f"Tamamlandı: {pptx_path}")


if __name__ == "__main__":
    input_pdf = "fransa.pdf"
    output_pptx = "fransa.pptx"

    pdf_to_pptx(input_pdf, output_pptx, dpi=200)